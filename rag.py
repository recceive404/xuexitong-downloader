"""
RAG 模块 — 课件知识库构建 + 问答

流程：
  ① 解析 PDF/PPT/TXT → 纯文本
  ② 按段落+语义边界分块
  ③ sentence-transformers 本地向量化（免费，无需 API）
  ④ 存入 ChromaDB（本地持久化）
  ⑤ 用户提问 → 检索 Top K 块 → DeepSeek V4 生成回答
"""

import os
import json
import warnings
from pathlib import Path

# 关闭 PyTorch pin_memory 警告（CPU 模式下无用）
warnings.filterwarnings("ignore", message=".*pin_memory.*")
warnings.filterwarnings("ignore", message=".*no accelerator.*")

# 国内 HuggingFace 镜像
if "HF_ENDPOINT" not in os.environ:
    os.environ["HF_ENDPOINT"] = "https://hf-mirror.com"

COURSES_DIR = Path(__file__).parent / "courses"
CHROMA_DIR = Path(__file__).parent / "chroma_db"


# ── 配置 ──────────────────────────────────────────────
CHUNK_SIZE = 500          # 每个文本块的最大字符数
CHUNK_OVERLAP = 80        # 相邻块的重叠字符数
TOP_K = 5                 # 检索返回的最大块数
EMBEDDING_MODEL = "paraphrase-multilingual-MiniLM-L12-v2"  # 本地模型，~120MB


class RAGEngine:
    def __init__(self):
        self.collection = None
        self.embedding_fn = None
        self._initialized = False

    def _ensure_initialized(self):
        if self._initialized:
            return

        # 延迟导入（这些库较重）
        from chromadb import PersistentClient
        from chromadb.utils import embedding_functions
        import sentence_transformers

        # 向量化函数 — 本地模型
        self.embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name=EMBEDDING_MODEL
        )

        # ChromaDB 客户端
        CHROMA_DIR.mkdir(parents=True, exist_ok=True)
        self.client = PersistentClient(path=str(CHROMA_DIR))

        self._initialized = True

    # ── 知识库构建 ─────────────────────────────────────

    def build_or_update(self):
        """扫描 courses/ 目录，构建或增量更新知识库"""
        self._ensure_initialized()

        if not COURSES_DIR.exists():
            print("[WARN] courses/ 目录不存在，请先下载课件")
            return

        # 获取所有可解析的文件
        all_files = self._scan_files()
        if not all_files:
            print("[WARN] 未找到任何课件文件")
            return

        print(f"[SCAN] 找到 {len(all_files)} 个课件文件")

        # 获取或创建 collection
        try:
            self.collection = self.client.get_collection(
                name="xuexitong_courses",
                embedding_function=self.embedding_fn
            )
            existing_count = self.collection.count()
            print(f"[DB] 现有知识库：{existing_count} 个文本块")
        except Exception:
            self.collection = self.client.create_collection(
                name="xuexitong_courses",
                embedding_function=self.embedding_fn
            )
            print("[DB] 创建新知识库")

        # 解析并入库
        total_chunks = 0
        for filepath in all_files:
            print(f"  [FILE] 解析：{filepath.relative_to(COURSES_DIR)}")

            text = self._parse_file(filepath)
            if not text or len(text) < 20:
                print(f"    [WARN] 文本内容过少，跳过")
                continue

            chunks = self._split_text(text, Path(filepath).stem)
            if not chunks:
                continue

            # 去重：已入库的文件直接跳过
            try:
                existing = self.collection.get(
                    where={"source": str(filepath.relative_to(COURSES_DIR))}
                )
                if existing and len(existing["ids"]) > 0:
                    print(f"    [SKIP] 已入库，跳过")
                    continue
            except Exception:
                pass

            # 批量添加
            try:
                self.collection.add(
                    ids=chunks["ids"],
                    documents=chunks["documents"],
                    metadatas=chunks["metadatas"]
                )
                total_chunks += len(chunks["ids"])
            except Exception as e:
                print(f"    [ERR] 入库失败：{e}")

        print(f"\n[OK] 知识库更新完成，共 {self.collection.count()} 个文本块")

    def _scan_files(self) -> list[Path]:
        """扫描 courses/ 下所有可解析文件。PDF 若缺 .txt 则先 OCR"""
        files = []
        extensions = {"*.txt", "*.pptx", "*.ppt", "*.html", "*.htm"}
        for ext in extensions:
            files.extend(COURSES_DIR.rglob(ext))

        # PDF：如果同名的 .txt 不存在，提取图片做 OCR
        for pdf_path in sorted(COURSES_DIR.rglob("*.pdf")):
            txt_path = pdf_path.with_suffix(".txt")
            if not txt_path.exists():
                print(f"  [OCR] PDF 缺 OCR：{pdf_path.name}，提取图片识别...")
                self._ocr_pdf(pdf_path, txt_path)
            if txt_path.exists():
                files.append(txt_path)
            else:
                # OCR 失败，尝试直接解析 PDF 文字
                files.append(pdf_path)

        return sorted(files)

    def _ocr_pdf(self, pdf_path: Path, txt_path: Path):
        """从 PDF 提取图片并 OCR，保存为 .txt"""
        try:
            import fitz  # PyMuPDF
            import easyocr
            from PIL import Image
            import io

            ocr = easyocr.Reader(['ch_sim', 'en'], gpu=False, verbose=False)
            doc = fitz.open(str(pdf_path))
            all_text = []

            for page_num in range(len(doc)):
                page = doc[page_num]
                # 先尝试直接提取文字
                text = page.get_text().strip()
                if len(text) > 50:
                    all_text.append(f"[第{page_num+1}页]\n{text}")
                    continue

                # 文字不够 → 从页面渲染图片做 OCR
                pix = page.get_pixmap(dpi=200)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                # 保存临时图片
                import tempfile
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    img.save(tmp.name)
                    try:
                        results = ocr.readtext(tmp.name, detail=0)
                        if results:
                            all_text.append(f"[第{page_num+1}页]\n" + "\n".join(results))
                    finally:
                        try:
                            Path(tmp.name).unlink()
                        except:
                            pass

                if (page_num + 1) % 5 == 0:
                    print(f"     OCR: {page_num+1}/{len(doc)} 页")

            doc.close()

            if all_text:
                txt_path.write_text("\n\n".join(all_text), encoding="utf-8")
                print(f"  [OK] {txt_path.name} ({len(all_text)} 页文字)")
            else:
                print(f"  [WARN] {pdf_path.name} 未识别到文字")

        except Exception as e:
            print(f"  [WARN] OCR 失败 {pdf_path.name}: {e}")

    def _parse_file(self, filepath: Path) -> str:
        """根据文件类型，解析为纯文本"""
        suffix = filepath.suffix.lower()

        if suffix == ".pdf":
            return self._parse_pdf(filepath)
        elif suffix in (".pptx", ".ppt"):
            return self._parse_pptx(filepath)
        else:
            # txt, html — 直接读
            return filepath.read_text(encoding="utf-8", errors="ignore")

    def _parse_pdf(self, filepath: Path) -> str:
        """PyMuPDF 解析 PDF"""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            print("    [ERR] 请安装 PyMuPDF：pip install PyMuPDF")
            return ""

        try:
            doc = fitz.open(str(filepath))
            text_parts = []
            for page in doc:
                text = page.get_text()
                if text:
                    text_parts.append(text)
            doc.close()
            return "\n\n".join(text_parts)
        except Exception as e:
            print(f"    [WARN] PDF 解析失败：{e}")
            return ""

    def _parse_pptx(self, filepath: Path) -> str:
        """python-pptx 解析 PPT"""
        try:
            from pptx import Presentation
        except ImportError:
            print("    [ERR] 请安装 python-pptx：pip install python-pptx")
            return ""

        try:
            prs = Presentation(str(filepath))
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
                if slide_texts:
                    text_parts.append(f"[幻灯片 {slide_num}]\n" + "\n".join(slide_texts))
            return "\n\n".join(text_parts)
        except Exception as e:
            print(f"    [WARN] PPT 解析失败：{e}")
            return ""

    def _split_text(self, text: str, source_name: str) -> dict | None:
        """将文本切分为重叠块"""
        try:
            from langchain_text_splitters import RecursiveCharacterTextSplitter
        except ImportError:
            from langchain.text_splitter import RecursiveCharacterTextSplitter

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=["\n\n", "\n", "。", ".", "！", "？", "，", ",", " ", ""],
            length_function=len,
        )

        chunks = splitter.split_text(text)
        if not chunks:
            return None

        ids = [f"{source_name}_{i}" for i in range(len(chunks))]
        metadatas = [{"source": source_name, "chunk_index": i} for i in range(len(chunks))]

        return {"ids": ids, "documents": chunks, "metadatas": metadatas}

    # ── 问答 ──────────────────────────────────────────

    def ask(self, question: str) -> str:
        """基于知识库检索 + DeepSeek V4 生成回答"""
        self._ensure_initialized()

        if not self.collection or self.collection.count() == 0:
            return "知识库为空，请先运行 build-rag 消化课件。"

        # 检索
        results = self.collection.query(
            query_texts=[question],
            n_results=TOP_K
        )

        docs = results.get("documents", [[]])[0]
        if not docs:
            return "未找到与问题相关的课件内容。"

        context = "\n\n---\n\n".join(docs)

        # 调用 DeepSeek V4 生成回答
        answer = self._generate_answer(question, context)
        return answer

    def _get_api_config(self):
        """获取 API 配置（环境变量 → ~/.claude/settings.json）"""
        api_url = os.environ.get("ANTHROPIC_BASE_URL")
        api_key = os.environ.get("ANTHROPIC_AUTH_TOKEN")
        model = os.environ.get("ANTHROPIC_MODEL")

        if not api_key:
            # 尝试从 Claude Code 配置读取
            settings_path = Path.home() / ".claude" / "settings.json"
            if settings_path.exists():
                try:
                    settings = json.loads(settings_path.read_text(encoding="utf-8"))
                    env = settings.get("env", {})
                    if not api_url:
                        api_url = env.get("ANTHROPIC_BASE_URL")
                    if not api_key:
                        api_key = env.get("ANTHROPIC_AUTH_TOKEN")
                    if not model:
                        model = env.get("ANTHROPIC_MODEL")
                except Exception:
                    pass

        return (
            api_url or "https://api.deepseek.com/anthropic",
            api_key or "",
            model or "deepseek-v4-pro[1m]"
        )

    def _generate_answer(self, question: str, context: str) -> str:
        """调用 DeepSeek V4（通过 Anthropic 兼容 API）"""
        print("[V3-NEW-CODE] _generate_answer called")
        import requests

        api_url, api_key, model = self._get_api_config()

        if not api_key:
            # 最后尝试：直接读文件
            try:
                from pathlib import Path as _Path
                settings = json.loads((_Path.home() / ".claude" / "settings.json").read_text(encoding="utf-8"))
                api_key = settings.get("env", {}).get("ANTHROPIC_AUTH_TOKEN", "")
                api_url = settings.get("env", {}).get("ANTHROPIC_BASE_URL", api_url)
            except:
                pass

        if not api_key:
            return (
                "[DEBUG] api_key为空，检查settings.json失败。\n\n"
                "以检索到的相关内容作为参考：\n\n"
                + context[:1500]
            )

        prompt = f"""你是一个学习助手，根据以下课件内容回答用户的问题。

课件内容：
{context}

用户问题：{question}

请基于课件内容回答。如果课件内容不足以回答，请明确说明。
用中文回答，简洁清晰。"""

        try:
            resp = requests.post(
                f"{api_url}/v1/messages",
                headers={
                    "x-api-key": api_key,
                    "anthropic-version": "2023-06-01",
                    "Content-Type": "application/json",
                },
                json={
                    "model": model,
                    "max_tokens": 1024,
                    "messages": [
                        {"role": "user", "content": prompt}
                    ]
                },
                timeout=60
            )
            if resp.status_code == 200:
                data = resp.json()
                return data["content"][0]["text"]
            else:
                # 降级：返回检索结果
                return (
                    f"[WARN] API 调用失败 ({resp.status_code})。\n\n"
                    f"以下是与问题相关的课件片段：\n\n"
                    f"{context[:2000]}"
                )
        except Exception as e:
            return (
                f"[WARN] API 调用异常：{e}\n\n"
                f"以下是与问题相关的课件片段：\n\n"
                f"{context[:2000]}"
            )

    # ── 工具方法 ──────────────────────────────────────

    def list_courses(self) -> list[str]:
        """列出已消化的课程（从 courses/ 目录中读取）"""
        if not COURSES_DIR.exists():
            return []
        return [
            d.name for d in COURSES_DIR.iterdir()
            if d.is_dir() and any(d.iterdir())
        ]

    def load(self):
        """加载已有知识库（用于问答，不需要重建）"""
        self._ensure_initialized()
        try:
            self.collection = self.client.get_collection(
                name="xuexitong_courses",
                embedding_function=self.embedding_fn
            )
            print(f"[OK] 知识库已加载（{self.collection.count()} 个文本块）")
        except Exception:
            self.collection = None
            print("[WARN] 知识库尚未构建，请先运行 build-rag")

    def persist(self):
        """ChromaDB 自动持久化，此方法仅做显式提示"""
        if self.collection:
            print(f"[SAVE] 知识库已持久化到 {CHROMA_DIR}（{self.collection.count()} 个块）")
