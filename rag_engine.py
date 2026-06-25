"""
轻量 RAG 模块 — 纯文本搜索 + DeepSeek Chat API

流程：
  ① 扫描 courses/ 下所有 .txt 文件 → 建立文本索引
  ② 用户提问 → 关键词匹配搜索 → 取最相关文本
  ③ 上下文 + 问题 → DeepSeek API → 生成回答

无需 ChromaDB / sentence-transformers / PyTorch / EasyOCR
"""
import os
import json
import re
from pathlib import Path

COURSES_DIR = Path(__file__).parent / "courses"
INDEX_FILE = Path(__file__).parent / "course_index.json"


class RAGEngine:
    """轻量级 RAG 引擎：文件索引 + 关键词搜索 + API 问答"""

    def __init__(self):
        self.documents: dict[str, str] = {}  # {文件名: 全文}
        self._initialized = False

    # ── 知识库构建 ─────────────────────────────────────

    def build_or_update(self):
        """扫描 courses/ 下所有可解析文件，建立文本索引

        支持格式：.txt .html .pdf .pptx .docx
        用法：把任意课件丢进 courses/ 下，运行 build-rag 即可
        """
        if not COURSES_DIR.exists():
            print("[WARN] courses/ 目录不存在，请先下载课件")
            return

        # 扫描所有支持的文件类型
        all_files = []
        for ext in ["*.txt", "*.html", "*.htm", "*.pdf", "*.pptx", "*.docx"]:
            all_files.extend(COURSES_DIR.rglob(ext))

        if not all_files:
            print("[WARN] 未找到任何课件文件（支持 txt/html/pdf/pptx/docx）")
            return

        print(f"[SCAN] 找到 {len(all_files)} 个文件")

        new_docs = {}
        for fp in sorted(all_files):
            key = str(fp.relative_to(COURSES_DIR)).replace("\\", "/")
            suffix = fp.suffix.lower()

            try:
                if suffix in (".txt", ".html", ".htm"):
                    text = fp.read_text(encoding="utf-8", errors="ignore")
                elif suffix == ".pdf":
                    text = self._parse_pdf(fp)
                elif suffix == ".pptx":
                    text = self._parse_pptx(fp)
                elif suffix == ".docx":
                    text = self._parse_docx(fp)
                else:
                    continue

                if text and len(text.strip()) > 10:
                    new_docs[key] = text
                    print(f"  [OK] {key} ({len(text):,} 字)")
                else:
                    print(f"  [LOW] {key} — 文本内容极少，可能以图片为主")
                    # 仍然保存，但标记为低质量
                    new_docs[key] = text or f"[此文件文本内容极少，建议查看原文件: {key}]"

            except Exception as e:
                print(f"  [ERR] {key}: {e}")
                new_docs[key] = f"[解析失败: {key}]"

        self.documents = new_docs
        self._initialized = True

        # 保存索引
        INDEX_FILE.write_text(
            json.dumps(self.documents, ensure_ascii=False), encoding="utf-8"
        )
        total_chars = sum(len(v) for v in self.documents.values())
        low_count = sum(1 for v in self.documents.values() if len(v) < 200)
        print(f"[OK] 索引完成：{len(self.documents)} 个文件，共 {total_chars:,} 字")
        if low_count:
            print(f"     {low_count} 个文件文本较少，提问时可能无法深度回答")

    def _parse_pdf(self, fp) -> str:
        """PyMuPDF 提取 PDF 文字"""
        try:
            import fitz
            doc = fitz.open(str(fp))
            parts = []
            for page in doc:
                t = page.get_text()
                if t:
                    parts.append(t)
            doc.close()
            return "\n".join(parts)
        except Exception as e:
            return f"[PDF解析失败: {e}]"

    def _parse_pptx(self, fp) -> str:
        """python-pptx 提取 PPT 文字"""
        try:
            from pptx import Presentation
            prs = Presentation(str(fp))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
                if slide_texts:
                    parts.append(f"[幻灯片 {i}]\n" + "\n".join(slide_texts))
            return "\n\n".join(parts)
        except Exception as e:
            return f"[PPT解析失败: {e}]"

    def _parse_docx(self, fp) -> str:
        """python-docx 提取 Word 文字"""
        try:
            from docx import Document
            doc = Document(str(fp))
            parts = []
            for para in doc.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
            # 也提取表格中的文字
            for table in doc.tables:
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        parts.append(" | ".join(row_texts))
            return "\n".join(parts)
        except Exception as e:
            return f"[DOCX解析失败: {e}]"

    def load(self):
        """加载已有索引"""
        if INDEX_FILE.exists():
            try:
                self.documents = json.loads(
                    INDEX_FILE.read_text(encoding="utf-8")
                )
                self._initialized = True
                total_chars = sum(len(v) for v in self.documents.values())
                print(f"[OK] 已加载索引：{len(self.documents)} 个文件，{total_chars:,} 字")
            except Exception:
                print("[WARN] 索引文件损坏，请重新 build-rag")
        else:
            print("[WARN] 尚未构建索引，请先运行 build-rag")

    # ── 问答 ──────────────────────────────────────────

    def ask(self, question: str) -> str:
        """搜索 + 问答"""
        if not self._initialized or not self.documents:
            return "知识库为空，请先运行 build-rag 消化课件。"

        # 1. 关键词搜索
        ranked = self._search(question)
        if not ranked:
            return "未找到与问题相关的课件内容。"

        # 2. 构建上下文
        context_parts = []
        low_text_files = []
        total_len = 0
        max_context = 60000  # 约 3 万 token，DeepSeek V4 有 100 万 token 上下文，绰绰有余
        for filename, score in ranked[:10]:
            text = self.documents[filename]
            snippet = self._extract_relevant(text, question, max(max_context - total_len, 1000))
            context_parts.append(f"【{filename}】\n{snippet}")
            total_len += len(snippet)
            # 标记低文本文件
            if len(text) < 200:
                low_text_files.append(filename)
            if total_len >= max_context:
                break

        context = "\n\n---\n\n".join(context_parts)
        print(f"  [SEARCH] 匹配 {len(ranked)} 个文件，上下文 {total_len:,} 字")

        # 3. 调用 API
        answer = self._generate_answer(question, context)

        # 4. 如果匹配中包含低文本文件，追加提示
        if low_text_files:
            hint = "\n\n---\n📎 以下文件文本内容较少（可能包含图片/表格/思维导图），建议查看原文件：\n"
            for f in low_text_files[:5]:
                hint += f"  • courses/{f}\n"
            answer += hint

        return answer

    def _search(self, question: str) -> list[tuple[str, int]]:
        """关键词匹配搜索，返回 [(文件名, 匹配分数), ...]

        中文用 2-4 字滑动窗口切词（n-gram），英文用空格分词。
        这样"番茄大棚栽培"会被切成 番茄/茄大/大棚/棚栽/栽培 等，
        每个片段都能独立匹配，不需要分词库。
        """
        keywords = []

        # 中文 n-gram：2字、3字、4字滑动窗口
        cn_chars = ''.join(c for c in question if '一' <= c <= '鿿')
        for n in (4, 3, 2):
            for i in range(len(cn_chars) - n + 1):
                keywords.append(cn_chars[i:i+n])

        # 英文/数字分词
        en_part = re.sub(r'[一-鿿]+', ' ', question)
        en_words = re.findall(r'[a-zA-Z0-9]{2,}', en_part)
        keywords.extend(en_words)

        # 去重+优先级：长词权重高
        seen = set()
        unique_kw = []
        for kw in keywords:
            if kw not in seen:
                seen.add(kw)
                unique_kw.append(kw)

        if not unique_kw:
            unique_kw = [question]

        # 搜索评分
        scored = []
        for filename, text in self.documents.items():
            score = 0
            text_lower = text.lower()
            fl_lower = filename.lower()
            for kw in unique_kw:
                kw_lower = kw.lower()
                count = text_lower.count(kw_lower)
                if count > 0:
                    # 长词匹配权重更高
                    weight = len(kw) * len(kw)
                    score += weight + count
                # 文件名匹配额外加分
                if kw_lower in fl_lower:
                    score += 15
            if score > 0:
                scored.append((filename, score))

        scored.sort(key=lambda x: x[1], reverse=True)
        return scored

    def _extract_relevant(self, text: str, question: str, max_len: int) -> str:
        """从文本中提取与问题最相关的片段（用 n-gram 定位）"""
        if len(text) <= max_len:
            return text

        # 用 2-gram 找最佳匹配位置
        cn_chars = ''.join(c for c in question if '一' <= c <= '鿿')
        best_pos = 0
        best_score = 0
        for n in (3, 2):
            for i in range(len(cn_chars) - n + 1):
                kw = cn_chars[i:i+n]
                pos = text.find(kw)
                if pos >= 0:
                    # 统计出现次数作为分数
                    count = text.count(kw)
                    score = count * n
                    if score > best_score:
                        best_score = score
                        best_pos = pos

        start = max(0, best_pos - max_len // 4)
        end = min(len(text), start + max_len)
        snippet = text[start:end]

        if start > 0:
            snippet = "…" + snippet
        if end < len(text):
            snippet += "…"

        return snippet

    # ── API 调用 ──────────────────────────────────────

    def _get_api_config(self):
        """获取 API 配置（.env → 环境变量 → ~/.claude/settings.json）"""
        try:
            from dotenv import load_dotenv
            load_dotenv(Path(__file__).parent / ".env")
        except Exception:
            pass

        api_url = os.environ.get("ANTHROPIC_BASE_URL")
        api_key = os.environ.get("ANTHROPIC_AUTH_TOKEN")
        model = os.environ.get("ANTHROPIC_MODEL")

        if not api_key:
            settings_path = Path.home() / ".claude" / "settings.json"
            if settings_path.exists():
                try:
                    settings = json.loads(settings_path.read_text(encoding="utf-8"))
                    env = settings.get("env", {})
                    api_url = api_url or env.get("ANTHROPIC_BASE_URL")
                    api_key = env.get("ANTHROPIC_AUTH_TOKEN")
                    model = model or env.get("ANTHROPIC_MODEL")
                except Exception:
                    pass

        return (
            api_url or "https://api.deepseek.com/anthropic",
            api_key or "",
            model or "deepseek-v4-pro[1m]"
        )

    def _generate_answer(self, question: str, context: str) -> str:
        """调用 DeepSeek V4"""
        import requests

        api_url, api_key, model = self._get_api_config()

        if not api_key:
            return (
                "⚠️ 未配置 API Key。\n\n"
                "以下是与问题相关的课件片段：\n\n"
                + context[:2000]
            )

        prompt = f"""你是一个学习助手，根据以下课件内容回答用户的问题。

课件内容：
{context}

用户问题：{question}

请基于课件内容回答。如果课件内容不足以回答，请明确说明。
用中文回答，简洁清晰。"""

        try:
            resp = requests.post(
                f"{api_url}/v1/messages",
                headers={
                    "x-api-key": api_key,
                    "anthropic-version": "2023-06-01",
                    "Content-Type": "application/json",
                },
                json={
                    "model": model,
                    "max_tokens": 1024,
                    "messages": [{"role": "user", "content": prompt}]
                },
                timeout=60
            )
            if resp.status_code == 200:
                return self._parse_response(resp.text)
            else:
                return (
                    f"[API 错误 {resp.status_code}]\n\n"
                    f"以下是与问题相关的课件片段：\n\n{context[:2000]}"
                )
        except Exception as e:
            return f"[API 异常：{e}]\n\n课件片段：\n{context[:2000]}"

    def _parse_response(self, raw: str) -> str:
        """解析 DeepSeek API 返回（兼容多种格式）

        DeepSeek V4 返回格式：[{type:'thinking',...}, {type:'text',...}]
        需要跳过 thinking 块，提取 text 块。
        """
        # 尝试 JSON
        try:
            data = json.loads(raw)
        except json.JSONDecodeError:
            try:
                import ast
                data = ast.literal_eval(raw)
            except Exception:
                return raw[:2000]

        # DeepSeek 列表格式：[{type:'thinking',...}, {type:'text',...}]
        if isinstance(data, list):
            # 优先找 text 块，跳过 thinking
            for item in data:
                if isinstance(item, dict) and item.get("type") == "text":
                    return item.get("text", "")
            # 没有 text 块，返回最后一个 thinking 摘要
            for item in reversed(data):
                if isinstance(item, dict) and item.get("type") == "thinking":
                    thinking = item.get("thinking", "")
                    if thinking:
                        return f"[思考过程]\n{thinking[:500]}"
            return str(data)[:500]

        # OpenAI 格式
        if isinstance(data, dict):
            if "choices" in data:
                return data["choices"][0]["message"]["content"]
            # Anthropic 格式
            content = data.get("content", "")
            if isinstance(content, list) and content:
                for item in content:
                    if isinstance(item, dict) and item.get("type") == "text":
                        return item.get("text", "")
                return content[0].get("text", str(content))
            if isinstance(content, str):
                return content

        return raw[:2000]

    # ── 工具方法 ──────────────────────────────────────

    def list_courses(self) -> list[str]:
        """列出已下载的课程"""
        if not COURSES_DIR.exists():
            return []
        return [
            d.name for d in COURSES_DIR.iterdir()
            if d.is_dir() and any(d.iterdir())
        ]

    def persist(self):
        """显式保存索引"""
        if self.documents:
            INDEX_FILE.write_text(
                json.dumps(self.documents, ensure_ascii=False),
                encoding="utf-8"
            )
            print(f"[SAVE] 索引已保存（{len(self.documents)} 个文件）")
